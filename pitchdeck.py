from pptx import Presentation
import csv

# Open the CSV file

# Print the resulting list

def replace_text_in_presentation(presentation_path, data, name):
    presentation = Presentation(presentation_path)
    key = list(data)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key in data:
                            if key in run.text:
                                run.text = run.text.replace(key, data[key])

    presentation.save("output/" + name+ " Pitch Deck.pptx")


def read_csv(file_path):
    data_list = []
    with open(file_path, 'r') as file:
        reader = csv.reader(file)
        for row in reader:
            data_list.append(row[0])
    return data_list


# Replace "old_text" with "new_text" in the PowerPoint presentation
def process():
    namelist = read_csv("company.csv")
    data = {
        '{NAME}' : "Kim Shi Tong",
        '{PRONOUN}' : "Mr",
        '{PHONE}' : "81754200",
        '{MONTH}' : "June",
        '{YEAR}' : "2023"
    }

    print(namelist)
    print(list(data))
    for name in namelist:
        data['{COMPANY}'] = name
        replace_text_in_presentation("pitchdeck.pptx", data,name)

process()